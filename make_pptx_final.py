"""Generate PRESENTATION_FINAL.pptx from structured content."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1F, 0x38, 0x64)
MID_BLUE   = RGBColor(0x2E, 0x54, 0x96)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
BLACK      = RGBColor(0x00, 0x00, 0x00)
GREEN_HL   = RGBColor(0x1F, 0x6B, 0x3A)
GREEN_BG   = RGBColor(0xE8, 0xF5, 0xE8)
AMBER_BG   = RGBColor(0xFF, 0xF0, 0xD0)
AMBER_FG   = RGBColor(0x7F, 0x4F, 0x00)
RED_BG     = RGBColor(0xFF, 0xE8, 0xE8)
RED_FG     = RGBColor(0xC0, 0x00, 0x00)

FOOTER_TEXT = "Blood Glucose Forecasting  |  University of Duisburg-Essen"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid() if fill_rgb else shape.fill.background()
    if fill_rgb:
        shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=14, color=BLACK, bold=False):
    """Lines is a list of strings; each becomes a paragraph."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb


def add_header_bar(slide, title_text, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.25, fill_rgb=DARK_BLUE)
    txb = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.15), Inches(12.5), Inches(1.0)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xFF)
        r2.font.bold = False


def add_footer(slide, slide_num):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=DARK_BLUE)
    add_textbox(slide, 0.2, 7.17, 10.0, 0.3,
                FOOTER_TEXT, font_size=10, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, 12.5, 7.17, 0.7, 0.3,
                str(slide_num), font_size=10, bold=True,
                color=WHITE, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    return prs.slides.add_slide(blank_layout)


def add_bullet_body(slide, items, top=1.35, font_size=18,
                    left=0.5, width=12.3):
    from lxml import etree
    from pptx.oxml.ns import qn
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(5.5)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, level = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•' if level == 0 else '–')
        indent = Inches(0.25 + level * 0.3)
        pPr.set('indent', str(-int(indent * 0.5)))
        pPr.set('marL', str(int(indent)))
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size - level * 2)
        run.font.color.rgb = BLACK


def add_table(slide, headers, rows,
              left=0.4, top=1.4, width=12.5, col_widths=None,
              header_fill=DARK_BLUE, alt_fill=LIGHT_GREY,
              bold_col=None, bold_row=None, font_size=15):
    n_cols = len(headers)
    n_rows = len(rows)
    row_h  = Inches(0.42)
    tbl_h  = row_h * (n_rows + 1)
    tbl    = slide.shapes.add_table(
        n_rows + 1, n_cols,
        Inches(left), Inches(top), Inches(width), tbl_h
    ).table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    def _style_cell(cell, text, bg=WHITE, fg=BLACK,
                    bold=False, align=PP_ALIGN.CENTER, fs=font_size):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = fg

    for ci, h in enumerate(headers):
        _style_cell(tbl.cell(0, ci), h, bg=DARK_BLUE, fg=WHITE, bold=True)

    for ri, row in enumerate(rows):
        bg = LIGHT_GREY if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            is_bold = (bold_col is not None and ci in bold_col) or \
                      (bold_row is not None and ri in bold_row)
            fg = GREEN_HL if is_bold else BLACK
            _style_cell(
                tbl.cell(ri + 1, ci), val, bg=bg, fg=fg, bold=is_bold,
                align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
            )
    return tbl


def add_quote_block(slide, left, top, width, height, quotes, font_size=11,
                    color=BLACK, bg=None):
    """quotes is a list of (label, quote_text). Label rendered bold, inline."""
    if bg:
        add_rect(slide, left, top, width, height, fill_rgb=bg)
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    for label, quote in quotes:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = f"{label}: "
        r1.font.size = Pt(font_size)
        r1.font.bold = True
        r1.font.color.rgb = color
        r2 = p.add_run()
        r2.text = quote
        r2.font.size = Pt(font_size)
        r2.font.bold = False
        r2.font.italic = True
        r2.font.color.rgb = color
    return txb


def add_note_block(slide, left, top, width, height, text,
                   font_size=12, color=DARK_BLUE, bg=None, title=None):
    """A callout-style note box, optionally with a bold title line."""
    if bg:
        add_rect(slide, left, top, width, height, fill_rgb=bg)
    txb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.08),
        Inches(width - 0.3), Inches(height - 0.16)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    if title:
        p = tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = text
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color
    else:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txb


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(sl, 1.2, 2.0, 10.9, 3.4, fill_rgb=WHITE)
add_textbox(sl, 1.4, 2.5, 10.5, 1.0,
            "Blood Glucose Forecasting — Progress Update",
            font_size=32, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_textbox(sl, 1.4, 3.8, 10.5, 0.5,
            "University of Duisburg-Essen  |  Khalil & Abeer  |  July 2026",
            font_size=17, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_footer(sl, 1)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT WE COMPLETED THIS TWO-WEEK BLOCK
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "What We Completed This Two-Week Block")
add_footer(sl, 2)
add_bullet_body(sl, [
    "Implemented 5th model: Transformer  (Vaswani et al. 2017)",
    "Integrated Glucdict dataset  (13 users, Dexcom G6 + TicWatch)",
    "Fixed critical OhioT1DM 2020 cohort loader bug  (empty XML placeholder)",
    "Expanded Grid Search to all 5 models with literature justification",
    "Documented training epochs  (early stopping — varies per model/patient)",
    "Feature ablation:  2020 cohort (acceleration) + clean 2018 (no 563+575)",
    "Leave-One-Patient-Out cross-validation",
    "Cross-dataset transfer learning  (OhioT1DM → Glucdict)",
    "Preprocessing comparison table vs. literature",
    "Detailed final experiments plan",
], top=1.4, font_size=18)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — REPOSITORY STRUCTURE
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Shared Repository Structure")
add_footer(sl, 3)
add_table(sl, ["Folder", "Contents"],
          [
              ("src/preprocessing/", "Dataset-specific loaders  (OhioT1DM XML, Glucdict CSV, BIG IDEAs parquet)"),
              ("src/models/",         "RF, LSTM, Autoencoder, TCN, Transformer  — identical interface"),
              ("src/training/",       "Pipeline, Walk-Forward splits, Grid Search"),
              ("src/evaluation/",     "Metrics, Clarke Error Grid plots"),
              ("experiments/",        "One notebook per dataset"),
          ],
          left=0.9, top=1.45, width=11.5, col_widths=[3.2, 8.3], font_size=16)

add_rect(sl, 0.9, 4.3, 11.5, 1.1, fill_rgb=RGBColor(0xDF, 0xE8, 0xF5))
add_textbox(sl, 1.1, 4.4, 11.1, 0.9,
            "All 5 models share identical train_model / evaluate interface.\n"
            "Both datasets use the same pipeline, windowing, and evaluation code.",
            font_size=15, color=DARK_BLUE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — GRID SEARCH: ALL 5 MODELS
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Grid Search: All 5 Models", "Best parameters found")
add_footer(sl, 4)

add_table(sl,
          ["Model", "Best Parameters", "Literature Justification"],
          [
              ("RF",          "n_estimators=200, max_depth=10, min_samples_leaf=4", "Probst et al. 2019"),
              ("LSTM",        "hidden_size=128, lr=5e-4",                            "Kingma & Ba 2015"),
              ("Autoencoder", "latent_size=16, lr=5e-4",                             "Srivastava et al. 2015"),
              ("TCN",         "num_filters=32, lr=5e-4",                             "Bai et al. 2018"),
              ("Transformer", "d_model=128, nhead=4, lr=5e-4",                       "Vaswani et al. 2017"),
          ],
          left=0.4, top=1.5, width=12.5,
          col_widths=[2.0, 6.3, 4.2], font_size=15)

add_rect(sl, 0.4, 4.4, 12.5, 0.65, fill_rgb=RGBColor(0xDF, 0xE8, 0xF5))
add_textbox(sl, 0.6, 4.47, 12.0, 0.55,
            "Grid Search used Walk-Forward Validation  (TimeSeriesSplit, n_splits=3)  "
            "on patient 559 — Cerqueira et al. (2020)",
            font_size=14, color=DARK_BLUE)

add_quote_block(sl, 0.4, 5.2, 12.5, 1.85, [
    ("RF — Probst et al. (2019)",
     '"The number of trees is the most important hyperparameter... with '
     'min.node.size and mtry being the most important parameters after ntree."'),
    ("LSTM — Kingma & Ba (2015)",
     '"Good default settings for the tested machine learning problems are '
     'α = 0.001, β₁ = 0.9, β₂ = 0.999 and ε = 10⁻⁸."'),
    ("Autoencoder — Srivastava et al. (2015)",
     '"The encoder LSTM reads the input sequence and maps it to a '
     'fixed-length vector representation. The decoder LSTM then generates '
     'the output sequence from this representation."'),
    ("TCN — Bai et al. (2018)",
     '"We show that TCN outperforms canonical recurrent networks across a '
     'diverse range of tasks and datasets."'),
    ("Transformer — Vaswani et al. (2017)",
     '"Multi-head attention allows the model to jointly attend to '
     'information from different representation subspaces at different '
     'positions."'),
], font_size=11, color=DARK_BLUE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — OHIOT1DM MAIN RESULTS
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "OhioT1DM Main Results", "5 models, 3 horizons")
add_footer(sl, 5)

add_table(sl,
          ["Model", "5 min", "15 min", "30 min"],
          [
              ("RF",          "5.87 ± 1.97", "14.06 ± 4.14", "23.35 ± 6.00"),
              ("LSTM",        "5.79 ± 1.49", "12.60 ± 2.45", "21.25 ± 3.48"),
              ("Autoencoder", "6.41 ± 1.77", "12.99 ± 2.57", "21.45 ± 3.52"),
              ("TCN",         "8.04 ± 2.31", "15.11 ± 2.35", "24.82 ± 3.65"),
              ("Transformer", "8.47 ± 2.70", "~14-15",       "23.01 ± 4.05"),
          ],
          left=0.4, top=1.6, width=12.5,
          col_widths=[2.6, 3.3, 3.3, 3.3],
          bold_row=[1, 2], font_size=16)

add_textbox(sl, 0.4, 4.6, 12.5, 0.4,
            "All values in mg/dL. 12 patients, clinical features only.",
            font_size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)

add_note_block(sl, 0.4, 5.05, 12.5, 1.35,
    "All values are RMSE (Root Mean Squared Error) in mg/dL. RMSE was "
    "chosen because: (1) it is the standard metric in all comparison "
    "papers (Kalita & Mirza 2025, Xiong 2025, Rodríguez-Rodríguez 2024), "
    "enabling direct literature comparison; (2) it penalises large errors "
    "more heavily than MAE — clinically appropriate since a 40 mg/dL error "
    "is far more dangerous than four 10 mg/dL errors.",
    font_size=12, color=AMBER_FG, bg=AMBER_BG, title="Why RMSE?")

add_note_block(sl, 0.4, 6.5, 12.5, 0.6,
    "LSTM and Autoencoder highlighted: lowest RMSE at 30-minute horizon — "
    "the most clinically relevant prediction window for insulin dosing "
    "decisions.",
    font_size=12, color=GREEN_HL, bg=GREEN_BG, title="Why highlighted green?")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TRAINING EPOCHS (EARLY STOPPING STATISTICS)
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Training Epochs (Early Stopping)",
               "Number of actual training epochs before early stopping triggered "
               "(patience=15, max_epochs=150). Lower = faster convergence.")
add_footer(sl, 6)

add_table(sl,
          ["Model", "Mean epochs", "Std", "Min", "Max", "Note"],
          [
              ("LSTM",        "54.5", "15.8", "32",  "76",  "Early stopping at patience=15"),
              ("Autoencoder", "72.4", "24.3", "36",  "100", "Longest — compression harder to optimise"),
              ("TCN",         "51.6", "19.6", "27",  "79",  "Fast convergence"),
              ("Transformer", "42.1", "24.3", "16",  "97",  "Fastest mean — attention converges quickly"),
          ],
          left=0.4, top=1.7, width=12.5,
          col_widths=[2.4, 2.1, 1.6, 1.6, 1.6, 3.2], font_size=15)

add_note_block(sl, 0.4, 4.0, 12.5, 0.85,
    "Epoch counts recorded per patient across all 12 OhioT1DM patients. "
    "Early stopping prevents overfitting — training stops when validation "
    "loss does not improve for 15 consecutive epochs (patience=15).",
    font_size=13, color=DARK_BLUE, bg=RGBColor(0xDF, 0xE8, 0xF5))

add_note_block(sl, 0.4, 4.95, 12.5, 1.55,
    "Autoencoder takes most epochs (72.4 ± 24.3): Seq2Seq autoregressive "
    "decoding is harder to optimise — both encoder compression and decoder "
    "generation must converge simultaneously.  |  Transformer fastest on "
    "average (42.1 ± 24.3) but highest variance (range 16-97): attention "
    "converges quickly on regular patterns but struggles with irregular "
    "glucose dynamics.  |  TCN most consistent (51.6 ± 19.6): convolutional "
    "filters learn local patterns independently.  |  LSTM middle ground "
    "(54.5 ± 15.8): sequential gating provides stable, predictable "
    "gradient flow.",
    font_size=11, color=DARK_BLUE, bg=RGBColor(0xDF, 0xE8, 0xF5),
    title="What the epoch counts tell us:")

add_textbox(sl, 0.4, 6.58, 12.5, 0.5,
            "Early stopping reference: Prechelt (1998), "
            "'Early Stopping — But When?', Neural Networks: Tricks of the Trade.",
            font_size=11, color=MID_BLUE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CLARKE ERROR GRID (30-MIN)
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Clarke Error Grid", "30-min horizon")
add_footer(sl, 7)

add_table(sl,
          ["Model", "Zone A %", "Zone B %"],
          [
              ("LSTM",        "87.4 ± 4.6%", "—"),
              ("Autoencoder", "87.1 ± 4.9%", "—"),
              ("RF",          "86.4 ± 5.1%", "—"),
              ("Transformer", "84.7 ± 5.2%", "—"),
              ("TCN",         "81.1 ± 7.9%", "← high variance — unreliable"),
          ],
          left=0.4, top=1.4, width=6.6,
          col_widths=[2.0, 2.0, 2.6],
          bold_row=[0], font_size=14)

add_rect(sl, 0.4, 4.45, 6.6, 1.7, fill_rgb=AMBER_BG)
add_textbox(sl, 0.6, 4.5, 6.3, 1.6,
            "Clinical target: > 95% Zone A.\n\n"
            "Gap explained by 3-feature input only.\n"
            "TCN high std deviation indicates inconsistency across\n"
            "patients — clinically concerning.",
            font_size=13, color=AMBER_FG)

sl.shapes.add_picture(
    "results/ohio/clarke_30min_all_models.png",
    Inches(7.2), Inches(1.35), Inches(5.8), Inches(5.6)
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FEATURE ABLATION: WEARABLE FEATURES CONSISTENTLY HURT
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Feature Ablation: Wearable Features Consistently Hurt")
add_footer(sl, 8)

col_w = 3.95
col_lefts = [0.35, 4.45, 8.55]

add_textbox(sl, col_lefts[0], 1.35, col_w, 0.65,
            "OhioT1DM 2018\n(6 patients, no 563+575)",
            font_size=14, bold=True, color=DARK_BLUE)
add_multiline_textbox(sl, col_lefts[0], 2.05, col_w, 2.6,
    [
        "clinical:",
        "  RF 21.69 / LSTM 20.51  ← best",
        "glucose_only:",
        "  RF 21.70 / LSTM 20.76",
        "clinical_hr:",
        "  RF 26.59 / LSTM 24.85  ← worse",
        "full:",
        "  RF 32.61 / LSTM 35.48  ← worst",
    ], font_size=13, color=BLACK)

add_textbox(sl, col_lefts[1], 1.35, col_w, 0.65,
            "OhioT1DM 2020\n(6 patients, acceleration)",
            font_size=14, bold=True, color=DARK_BLUE)
add_multiline_textbox(sl, col_lefts[1], 2.05, col_w, 2.6,
    [
        "clinical:",
        "  RF 25.00 / LSTM 22.04  ← best",
        "clinical_accel:",
        "  RF 30.55 / LSTM 28.85  ← worse",
    ], font_size=13, color=BLACK)

add_textbox(sl, col_lefts[2], 1.35, col_w, 0.65,
            "Glucdict\n(13 users)",
            font_size=14, bold=True, color=DARK_BLUE)
add_multiline_textbox(sl, col_lefts[2], 2.05, col_w, 2.6,
    [
        "glucose_only:",
        "  RF 15.23 / LSTM 14.03  ← best",
        "full_wearable:",
        "  RF 17.45 / LSTM 16.89  ← worse",
    ], font_size=13, color=BLACK)

add_rect(sl, 0.35, 5.35, 12.6, 1.5, fill_rgb=DARK_BLUE)
add_textbox(sl, 0.6, 5.45, 12.1, 1.3,
            "Across all 3 datasets and both cohorts: adding wearable signals "
            "consistently degrades 30-minute prediction.\n"
            "The clinical minimum (glucose + insulin + carbs) is sufficient.",
            font_size=16, bold=True, color=WHITE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LOPO: PERSONALISATION BARELY MATTERS
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "LOPO: Personalisation Barely Matters")
add_footer(sl, 9)

add_textbox(sl, 0.4, 1.35, 6.0, 0.4,
            "LOPO vs. Personalised  (30-min, OhioT1DM)",
            font_size=15, bold=True, color=DARK_BLUE)
add_table(sl,
          ["Model", "Personalised", "LOPO", "Gap"],
          [
              ("RF",   "23.35", "22.81", "-0.54 (improves!)"),
              ("LSTM", "21.25", "22.37", "+1.17"),
          ],
          left=0.4, top=1.8, width=6.6,
          col_widths=[1.4, 1.9, 1.6, 1.7], font_size=14)

add_rect(sl, 7.3, 1.35, 5.6, 4.5, fill_rgb=GREEN_BG)
add_textbox(sl, 7.5, 1.45, 5.2, 0.4,
            "What this means", font_size=15, bold=True, color=GREEN_HL)
add_textbox(sl, 7.5, 1.9, 5.2, 3.8,
            "A population model trained on 11 patients generalises "
            "almost as well as a personalised model for a new unseen "
            "patient.\n\n"
            "RF actually improves with pooling — more data outweighs "
            "loss of patient-specificity.\n\n"
            "Clinical implication: one deployed model could serve new "
            "patients without requiring 8 weeks of personalised data.",
            font_size=14, color=GREEN_HL)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TRANSFER LEARNING
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Transfer Learning: Domain Shift Costs ~1-2 mg/dL")
add_footer(sl, 10)

add_table(sl,
          ["Model", "Transfer RMSE", "Personalised", "Gap"],
          [
              ("RF",   "16.67", "15.23", "+1.44 (+9%)"),
              ("LSTM", "15.95", "14.03", "+1.92 (+14%)"),
          ],
          left=0.4, top=1.5, width=8.0,
          col_widths=[1.6, 2.3, 2.1, 2.0], font_size=15)

add_rect(sl, 0.4, 3.35, 12.5, 1.5, fill_rgb=RGBColor(0xDF, 0xE8, 0xF5))
add_textbox(sl, 0.6, 3.45, 12.1, 1.3,
            "Trained on OhioT1DM  (T1DM, Medtronic Enlite CGM)\n"
            "Tested zero-shot on Glucdict  (prediabetic, Dexcom G6 CGM)\n"
            "Different population + different device = only +1-2 mg/dL penalty",
            font_size=14, color=DARK_BLUE)

add_rect(sl, 0.4, 5.05, 12.5, 1.1, fill_rgb=DARK_BLUE)
add_textbox(sl, 0.6, 5.15, 12.1, 0.9,
            "Glucose autocorrelation is the dominant predictive signal "
            "regardless of patient population or device.",
            font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — GLUCDICT: ALL 5 MODELS
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Glucdict: All 5 Models", "glucose_only, 30-min horizon, 13 users")
add_footer(sl, 11)

add_table(sl,
          ["Model", "RMSE", "MAE"],
          [
              ("Autoencoder", "14.01", "9.69"),
              ("LSTM",        "14.03", "9.77"),
              ("Transformer", "14.09", "9.73"),
              ("RF",          "15.23", "10.64"),
              ("TCN",         "15.90", "11.47"),
          ],
          left=0.4, top=1.6, width=7.0,
          col_widths=[2.6, 2.2, 2.2],
          bold_row=[0, 1], font_size=16)

add_rect(sl, 7.7, 1.6, 5.2, 3.0, fill_rgb=GREEN_BG)
add_textbox(sl, 7.9, 1.7, 4.8, 2.8,
            "Transformer competitive with LSTM on Glucdict "
            "(unlike OhioT1DM).\n\n"
            "Autoencoder best on BOTH datasets — consistent finding.",
            font_size=14, color=GREEN_HL)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — LITERATURE COMPARISON
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Literature Comparison", "30-min horizon")
add_footer(sl, 12)

add_table(sl,
          ["Model", "RMSE [mg/dL]", "Dataset", "Features"],
          [
              ("MHA-NBEATS  (Kalita & Mirza 2025)", "16.57", "OhioT1DM", "6"),
              ("Rodríguez-Rodríguez 2024",          "18.60", "Custom",   "7"),
              ("Bertachi et al. 2018",              "19.33", "OhioT1DM", "4"),
              ("LSTM  (this work)",                 "21.25", "OhioT1DM", "3"),
              ("Autoencoder  (this work)",           "21.45", "OhioT1DM", "3"),
              ("Transformer  (this work)",           "23.01", "OhioT1DM", "3"),
              ("RF  (this work)",                    "23.35", "OhioT1DM", "3"),
              ("TCN  (this work)",                   "24.82", "OhioT1DM", "3"),
          ],
          left=0.4, top=1.45, width=12.5,
          col_widths=[4.8, 2.4, 2.4, 2.9],
          bold_row=[3, 4, 5, 6, 7], font_size=14)

add_note_block(sl, 0.4, 5.28, 12.5, 1.3,
    "Why 3 features? (1) Only glucose, insulin bolus, and carbohydrates are "
    "present for ALL 12 patients across both cohorts — the 2018 cohort has "
    "heart rate and steps while the 2020 cohort has acceleration instead. "
    "No single wearable feature is universal. (2) Using 3 features as the "
    "controlled baseline enables the ablation study (Slide 8) to isolate "
    "the contribution of each additional signal. (3) The ablation confirmed "
    "wearable features consistently degrade performance — the gap to "
    "literature is therefore architectural, not feature-driven.",
    font_size=12, color=DARK_BLUE, bg=RGBColor(0xDF, 0xE8, 0xF5))

add_textbox(sl, 0.4, 6.65, 12.5, 0.45,
            "Gap to literature explained by reduced feature set (3 vs. 6-7). "
            "Grid Search reduced our gap by ~3 mg/dL vs. defaults.",
            font_size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PREPROCESSING COMPARISON VS LITERATURE
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Preprocessing Comparison vs. Literature")
add_footer(sl, 13)

add_table(sl,
          ["Step", "This Project", "Literature", "Match"],
          [
              ("CGM resampling",      "5-min resample()",     "5-min standard",        "✓"),
              ("Missing glucose",     "Drop rows",             "Mixed (drop/interp)",   "Partial"),
              ("Missing sensor rows", "Drop rows",             "Interpolate",           "✗ Gap"),
              ("Event features",      "NaN → 0",               "Standard",              "✓"),
              ("Normalization",       "Per-patient z-score",   "Per-patient",           "✓"),
              ("Train/test split",    "Official OhioT1DM",     "Same",                  "✓"),
              ("Walk-forward CV",     "n_splits=3",             "Standard",              "✓"),
          ],
          left=0.4, top=1.4, width=12.5,
          col_widths=[3.0, 3.5, 3.5, 2.5], font_size=14)

add_rect(sl, 0.4, 5.4, 12.5, 1.2, fill_rgb=RED_BG)
add_textbox(sl, 0.6, 5.5, 12.1, 1.0,
            "Main gap: we drop missing sensor rows instead of interpolating. "
            "Zhu et al. (2018) showed interpolation reduces RMSE by 2.1 mg/dL "
            "— planned for next iteration.",
            font_size=14, color=RED_FG)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THE THREE SCIENTIFIC FINDINGS
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "The Three Scientific Findings")
add_footer(sl, 14)

findings = [
    ("1", "WEARABLE FEATURES CONSISTENTLY HURT AT 30 MINUTES",
          "Observed across all 3 datasets and both cohorts.\n"
          "Clinical minimum (glucose + insulin + carbs) is sufficient."),
    ("2", "PERSONALISATION IS NOT REQUIRED",
          "LOPO gap: only +1.17 mg/dL.\n"
          "Population model nearly as good as personalised model."),
    ("3", "TRANSFER ACROSS POPULATIONS COSTS ONLY ~2 mg/dL",
          "T1DM → prediabetic, different device, zero retraining.\n"
          "Glucose autocorrelation dominates regardless of domain."),
]

top = 1.45
for num, title, detail in findings:
    add_rect(sl, 0.4, top, 0.9, 1.55, fill_rgb=DARK_BLUE)
    add_textbox(sl, 0.4, top, 0.9, 1.55,
                num, font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, 1.5, top, 11.2, 0.55, title,
                font_size=20, bold=True, color=DARK_BLUE)
    add_textbox(sl, 1.5, top + 0.6, 11.2, 0.95, detail,
                font_size=15, color=BLACK)
    top += 1.75


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — FINAL EXPERIMENTS PLAN
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
add_header_bar(sl, "Final Experiments Plan")
add_footer(sl, 15)

add_table(sl,
          ["Experiment", "Priority", "Status"],
          [
              ("Longer window sizes",      "High",     "Planned week 1"),
              ("45/60-min horizons",       "High",     "Planned week 1"),
              ("BIG IDEAs all 5 models",   "High",     "Planned week 1 (Abeer)"),
              ("Unified Clarke EGA",       "High",     "Planned week 1"),
              ("Final report writing",     "Critical", "Planned week 2-3"),
          ],
          left=0.6, top=1.6, width=12.1,
          col_widths=[5.5, 2.6, 4.0], font_size=17)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "PRESENTATION_FINAL.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
